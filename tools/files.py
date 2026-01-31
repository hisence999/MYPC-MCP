import os
import shutil
from datetime import datetime
from mcp.server.fastmcp import FastMCP

# Try to import Office libraries (fail silently if not installed)
try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# Default workspace directory (will be overridden by config)
DEFAULT_WORKSPACE = "D:\\ALICE"  # Fallback, will be replaced by config

# Default safe zones - can be configured
DEFAULT_SAFE_ZONES = [
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/Desktop"),
]


def register_file_tools(mcp: FastMCP, safe_zones: list[str] = None, base_url: str = "", default_workspace: str = None):
    """
    Register file management tools with tiered permissions.

    Permission Model:
    - READ operations: Allowed anywhere
    - WRITE operations: Only in safe zones
    - COPY: Can copy INTO safe zones from anywhere, but NOT out of safe zones

    Args:
        mcp: FastMCP instance
        safe_zones: List of directories where write operations are allowed.
        base_url: Base URL of the server (e.g., "http://localhost:9999")
        default_workspace: Default workspace directory for file operations
    """
    global DEFAULT_WORKSPACE

    allowed_zones = safe_zones or DEFAULT_SAFE_ZONES
    # Normalize paths
    allowed_zones = [os.path.abspath(z) for z in allowed_zones]

    # Use configured workspace if provided
    if default_workspace:
        DEFAULT_WORKSPACE = default_workspace

    # Ensure workspace exists
    os.makedirs(DEFAULT_WORKSPACE, exist_ok=True)

    def is_in_safe_zone(path: str) -> bool:
        """Check if path is within workspace or any safe zone. Workspace takes priority."""
        try:
            abs_path = os.path.abspath(path)
            # Normalize path case for Windows
            norm_path = os.path.normcase(abs_path)

            # Check workspace first (workspace is above safe zones)
            norm_workspace = os.path.normcase(os.path.abspath(DEFAULT_WORKSPACE))
            if norm_path == norm_workspace:
                return True
            if not norm_workspace.endswith(os.sep):
                norm_workspace += os.sep
            if norm_path.startswith(norm_workspace):
                return True

            # Then check safe zones
            for zone in allowed_zones:
                abs_zone = os.path.abspath(zone)
                norm_zone = os.path.normcase(abs_zone)

                # Check if it's the zone itself
                if norm_path == norm_zone:
                    return True

                # Check if it's a file inside the zone
                # Add separator to ensure we don't match partial folder names
                # e.g. C:\Users\User\Doc shouldn't match C:\Users\User\Documents
                if not norm_zone.endswith(os.sep):
                    norm_zone += os.sep

                if norm_path.startswith(norm_zone):
                    return True

            return False
        except Exception:
            return False

    def get_safe_zones_str() -> str:
        """Get formatted list of workspace and safe zones for error messages."""
        zones = [f"Workspace: {DEFAULT_WORKSPACE}"]
        zones.extend(f"  - {z}" for z in allowed_zones)
        return "\n".join(zones)

    # ==================== HELPERS FOR OFFICE FILES ====================

    def _read_docx(path: str, max_lines: int) -> str:
        if not docx:
            return "Error: python-docx not installed. Cannot read .docx files."
        try:
            doc = docx.Document(path)
            full_text = []
            for para in doc.paragraphs:
                if len(full_text) >= max_lines:
                    full_text.append(f"\n... (truncated at {max_lines} paragraphs)")
                    break
                full_text.append(para.text)
            return "\n".join(full_text)
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"

    def _read_xlsx(path: str, max_lines: int) -> str:
        if not openpyxl:
            return "Error: openpyxl not installed. Cannot read .xlsx files."
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            result = []

            # Read active sheet only
            sheet = wb.active
            result.append(f"[Sheet: {sheet.title}]")

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= max_lines:
                    result.append(f"\n... (truncated at {max_lines} rows)")
                    break

                # Filter None values and join
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                result.append(row_text)
                row_count += 1

            return "\n".join(result)
        except Exception as e:
            return f"Error reading XLSX: {str(e)}"

    def _read_pptx(path: str, max_lines: int) -> str:
        if not pptx:
            return "Error: python-pptx not installed. Cannot read .pptx files."
        try:
            prs = pptx.Presentation(path)
            text_runs = []

            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                text_runs.append(f"--- Slide {slide_count} ---")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)

                if len(text_runs) >= max_lines:
                    text_runs.append(f"\n... (truncated at {max_lines} text blocks)")
                    break

            return "\n".join(text_runs)
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"

    def _read_pdf(path: str, max_lines: int) -> str:
        if not PyPDF2:
            return "Error: PyPDF2 not installed. Cannot read .pdf files."
        try:
            reader = PyPDF2.PdfReader(path)
            text = []

            count = 0
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {count + 1} ---")
                    text.append(page_text)
                    count += 1

                if count >= max_lines: # Treat pages as lines/blocks
                    text.append(f"\n... (truncated at {max_lines} pages)")
                    break

            return "\n".join(text)
        except Exception as e:
            return f"Error reading PDF: {str(e)}"

    # ==================== WORKSPACE INFO ====================

    @mcp.tool(name="MyPC-get_workspace")
    def get_workspace() -> str:
        """
        Get the default workspace directory path.

        Default workspace is D:\ALICE. All file operations can use this as the base directory.
        You can use relative paths or just the workspace for operations.

        Returns:
            Default workspace path and information.
        """
        return f"Default Workspace: {DEFAULT_WORKSPACE}\n\nThis is the base directory for all file operations."

    # ==================== READ OPERATIONS (Anywhere) ====================

    @mcp.tool(name="MyPC-list_directory")
    def list_directory(path: str = None) -> str:
        """
        List contents of a directory. (READ - allowed anywhere)

        Default workspace: D:\ALICE (use path="D:\\ALICE" or omit for workspace)

        Args:
            path: Absolute path to the directory. If omitted or empty, uses workspace (D:\\ALICE).

        Returns:
            Formatted list of files and folders with their sizes and modification times.
        """
        # Use workspace if path not provided
        if not path:
            path = DEFAULT_WORKSPACE

        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        if not os.path.isdir(path):
            return f"Error: Not a directory: {path}"

        try:
            entries = []
            for name in os.listdir(path):
                full_path = os.path.join(path, name)
                try:
                    stat = os.stat(full_path)
                    mtime = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")

                    if os.path.isdir(full_path):
                        entries.append(f"[DIR]  {name}/  ({mtime})")
                    else:
                        size = stat.st_size
                        if size < 1024:
                            size_str = f"{size} B"
                        elif size < 1024 * 1024:
                            size_str = f"{size // 1024} KB"
                        else:
                            size_str = f"{size // (1024 * 1024)} MB"
                        entries.append(f"[FILE] {name}  ({size_str}, {mtime})")
                except PermissionError:
                    entries.append(f"[???]  {name}  (Permission Denied)")

            if not entries:
                return f"Directory is empty: {path}"

            return f"Contents of {path}:\n" + "\n".join(sorted(entries))

        except PermissionError:
            return f"Error: Permission denied to access: {path}"
        except Exception as e:
            return f"Error listing directory: {str(e)}"

    @mcp.tool(name="MyPC-read_file")
    def read_file(path: str, start_line: int = None, end_line: int = None, max_lines: int = 500) -> str:
        """
        Read contents of a file. Supports Text, Word (.docx), Excel (.xlsx), PPT (.pptx), PDF.
        (READ - allowed anywhere)

        Default workspace: D:\ALICE

        Args:
            path: Absolute path to the file. Can be in workspace (D:\\ALICE\\file.txt).
            start_line: Starting line number (1-indexed, None = from beginning).
            end_line: Ending line number (1-indexed, exclusive, None = read default max_lines lines).
            max_lines: Default maximum lines when start_line/end_line not specified (default 500).

        Examples:
            - (None, None): Read first max_lines lines (default 500)
            - (100, 200): Read lines 100-199
            - (100, None): Read max_lines lines starting from line 100
            - (None, 100): Read first 100 lines

        Returns:
            File contents as text.
        """
        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        if not os.path.isfile(path):
            return f"Error: Not a file: {path}"

        ext = os.path.splitext(path)[1].lower()

        # Default: read first max_lines lines
        default_max_lines = max_lines

        # Calculate actual line range
        if start_line is None and end_line is None:
            # Default: first max_lines lines
            read_start = 0
            read_end = default_max_lines
            range_info = f"Lines 1-{default_max_lines}"
        elif start_line is None and end_line is not None:
            # From beginning to end_line
            read_start = 0
            read_end = end_line
            range_info = f"Lines 1-{end_line}"
        elif start_line is not None and end_line is None:
            # From start_line for max_lines lines
            read_start = start_line - 1  # Convert to 0-indexed
            read_end = start_line - 1 + default_max_lines
            range_info = f"Lines {start_line}-{start_line + default_max_lines - 1}"
        else:
            # From start_line to end_line
            read_start = start_line - 1  # Convert to 0-indexed
            read_end = end_line
            range_info = f"Lines {start_line}-{end_line - 1}"

        # Route to appropriate handler (Office files don't support range reading)
        if ext == ".docx":
            return _read_docx(path, read_end)
        elif ext == ".xlsx":
            return _read_xlsx(path, read_end)
        elif ext == ".pptx":
            return _read_pptx(path, read_end)
        elif ext == ".pdf":
            return _read_pdf(path, read_end)

        # Default text reading for other files
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                lines = []

                # Skip to start line
                if read_start > 0:
                    for _ in range(read_start):
                        try:
                            next(f)
                        except StopIteration:
                            return f"Error: start_line {start_line} exceeds file length."

                # Read lines from start to end
                line_count = 0
                for line in f:
                    if line_count >= (read_end - read_start):
                        break
                    lines.append(line.rstrip('\n\r'))
                    line_count += 1

                if line_count == 0:
                    return f"Error: No lines found in specified range {range_info}."

            return f"{range_info}:\n" + "\n".join(lines)

        except PermissionError:
            return f"Error: Permission denied to read: {path}"
        except Exception as e:
            return f"Error reading file: {str(e)}"

    @mcp.tool(name="MyPC-get_file_info")
    def get_file_info(path: str) -> str:
        """
        Get detailed information about a file or directory. (READ - allowed anywhere)

        Default workspace: D:\ALICE

        Args:
            path: Path to the file or directory. Can be in workspace (D:\\ALICE\\file.txt).

        Returns:
            Detailed file information.
        """
        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        try:
            stat = os.stat(path)
            abs_path = os.path.abspath(path)
            in_safe = "Yes" if is_in_safe_zone(path) else "No"

            info = [
                f"Path: {abs_path}",
                f"Type: {'Directory' if os.path.isdir(path) else 'File'}",
                f"Size: {stat.st_size} bytes",
                f"Created: {datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Modified: {datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Accessed: {datetime.fromtimestamp(stat.st_atime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"In Safe Zone: {in_safe}",
            ]

            if os.path.isdir(path):
                try:
                    count = len(os.listdir(path))
                    info.append(f"Items: {count}")
                except PermissionError:
                    info.append("Items: (Permission Denied)")

            return "\n".join(info)

        except PermissionError:
            return f"Error: Permission denied to access: {path}"
        except Exception as e:
            return f"Error getting file info: {str(e)}"

    @mcp.tool(name="MyPC-list_safe_zones")
    def list_safe_zones() -> str:
        """
        List all configured safe zones where write operations are allowed.

        Returns:
            List of safe zone directories.
        """
        return "Safe Zones (write operations allowed):\n" + get_safe_zones_str()

    @mcp.tool(name="MyPC-get_download_url")
    def get_download_url(path: str) -> str:
        """
        Get a direct download URL for a file.

        The file MUST be located in a Safe Zone.

        Args:
            path: Absolute path to the file.

        Returns:
            Download URL if valid, or error message.
        """
        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        if not os.path.isfile(path):
            return f"Error: Not a file: {path}"

        if not is_in_safe_zone(path):
            return f"Error: Download denied. File must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        try:
            # URL encode the path
            import urllib.parse
            encoded_path = urllib.parse.quote(path)
            url = f"{base_url}/download?path={encoded_path}"
            return f"Download URL: {url}"
        except Exception as e:
            return f"Error generating URL: {str(e)}"

    # ==================== WRITE OPERATIONS (Safe Zones Only) ====================

    @mcp.tool(name="MyPC-edit_file")
    def edit_file(path: str, old_text: str, new_text: str, count: int = 1) -> str:
        """
        Edit a file by replacing specific text. (WRITE - safe zones only)

        Performs an exact string replacement. Preserves the rest of the file content.

        Args:
            path: Absolute path to the file.
            old_text: The text to be replaced (must match exactly, including newlines/indentation).
            new_text: The new text to insert.
            count: Number of occurrences to replace (default 1). Use -1 to replace all.

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(path):
            return f"Error: Edit operation denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        try:
            # Read file
            with open(path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Check if old_text exists
            if old_text not in content:
                # Provide a helpful hint if it's a whitespace issue
                normalized_content = " ".join(content.split())
                normalized_old = " ".join(old_text.split())
                if normalized_old in normalized_content:
                    return "Error: Text found but with different whitespace/indentation. Please provide exact match."
                return "Error: The 'old_text' was not found in the file."

            # Perform replacement
            new_content = content.replace(old_text, new_text, count)

            # Check if anything actually changed (in case count was 0 or other logic)
            if new_content == content:
                return "Warning: No changes were made (replacement logic returned same content)."

            # Write back
            with open(path, 'w', encoding='utf-8') as f:
                f.write(new_content)

            return f"File edited successfully: {path}"

        except Exception as e:
            return f"Error editing file: {str(e)}"

    @mcp.tool(name="MyPC-create")
    def create(items) -> str:
        """
        Create files and/or directories. (WRITE - safe zones only)

        IMPORTANT: This tool is for CREATING NEW FILES ONLY.
        - Use 'edit_file' to modify existing files
        - This tool will NOT modify existing files (fails if file exists)

        Default workspace: D:\ALICE (recommended for new items)
        Safe zones: Documents, Downloads, Desktop, Pictures, D:\\, E:\\

        ═══════════════════════════════════════════════════════════════
        📦 支持的格式 | Supported Formats
        ═══════════════════════════════════════════════════════════════

        1️⃣ 简单字符串（自动判断类型）| Simple string (auto-detect):
           - "folder"          → 创建目录 | Create directory
           - "file.txt"        → 创建空文件 | Create empty file
           - "D:\\path\\item"  → 使用绝对路径 | Use absolute path

        2️⃣ 文件内容字典 | File content dict:
           {"file.txt": "content", "config.json": '{"key": "value"}'}

        3️⃣ 列表格式（混合）| List format (mixed):
           [
               "folder1",                           # 目录 | Directory
               "folder2",
               {"file.txt": "content"},             # 文件 | File
               {"D:\\ALICE\\data.json": "{}"}       # 绝对路径文件 | Absolute path
           ]

        4️⃣ 详细格式（明确指定类型）| Detailed format (explicit type):
           [
               {"type": "dir", "path": "folder"},
               {"type": "file", "path": "test.txt", "content": "hello"}
           ]

        ═══════════════════════════════════════════════════════════════
        💡 示例 | Examples
        ═══════════════════════════════════════════════════════════════

        # 创建单个目录 | Create single directory
        create("projects")

        # 创建多个目录和文件 | Create multiple directories and files
        create([
            "src",
            "docs",
            {"README.md": "# My Project"},
            {"config.json": '{"debug": true}'}
        ])

        # 详细格式 | Detailed format
        create([
            {"type": "dir", "path": "src"},
            {"type": "file", "path": "index.js", "content": "console.log('hi')"}
        ])

        Returns:
            Success or error message with details for each item.
        """
        # Normalize input to list of items
        item_list = []
        results = []
        file_success = file_skip = file_fail = file_denied = 0
        dir_success = dir_fail = dir_denied = 0

        # Helper to process items recursively
        def process_item(item):
            nonlocal file_success, file_skip, file_fail, file_denied
            nonlocal dir_success, dir_fail, dir_denied

            # String: simple path (auto-detect type)
            if isinstance(item, str):
                # If no extension, treat as directory
                if '.' not in os.path.basename(item) or item.endswith('/'):
                    return {"type": "dir", "path": item}
                else:
                    return {"type": "file", "path": item, "content": ""}

            # Dict with explicit type
            elif isinstance(item, dict):
                # Check if it's the detailed format
                if "type" in item:
                    return item
                # Otherwise, it's a file content dict {"path": "content"}
                else:
                    items_list = []
                    for path, content in item.items():
                        items_list.append({"type": "file", "path": path, "content": content})
                    return items_list

            # List: process each element
            elif isinstance(item, list):
                items_list = []
                for sub in item:
                    processed = process_item(sub)
                    if isinstance(processed, list):
                        items_list.extend(processed)
                    else:
                        items_list.append(processed)
                return items_list

            return None

        # Process input
        processed = process_item(items)
        if isinstance(processed, list):
            item_list = processed
        elif processed:
            item_list = [processed]

        if not item_list:
            return "Error: No valid items specified."

        # Process each item
        for item in item_list:
            if not item or not isinstance(item, dict):
                continue

            item_type = item.get("type")
            path = item.get("path", "")
            content = item.get("content", "")

            if not path:
                results.append(f"[SKIP] Missing path")
                continue

            # If path is not absolute, use workspace
            if not os.path.isabs(path):
                path = os.path.join(DEFAULT_WORKSPACE, path)

            # Check safe zone
            if not is_in_safe_zone(path):
                results.append(f"[DENIED] Not in safe zone: {path}")
                if item_type == "file":
                    file_denied += 1
                else:
                    dir_denied += 1
                continue

            # Create directory
            if item_type == "dir":
                try:
                    os.makedirs(path, exist_ok=True)
                    results.append(f"[OK] Directory created: {path}")
                    dir_success += 1
                except Exception as e:
                    results.append(f"[FAIL] Directory {path}: {str(e)}")
                    dir_fail += 1

            # Create file
            elif item_type == "file":
                # Check if file already exists
                if os.path.exists(path):
                    results.append(f"[SKIP] File already exists: {path} (Use 'edit_file' to modify)")
                    file_skip += 1
                    continue

                try:
                    # Create parent directory if needed
                    parent = os.path.dirname(path)
                    if parent and not os.path.exists(parent):
                        os.makedirs(parent, exist_ok=True)

                    # Write file
                    with open(path, 'w', encoding='utf-8') as f:
                        f.write(content)

                    results.append(f"[OK] File created: {os.path.basename(path)} -> {path}")
                    file_success += 1

                except Exception as e:
                    results.append(f"[FAIL] File {os.path.basename(path)}: {str(e)}")
                    file_fail += 1

        # Summary
        file_summary = f"Files: {file_success} created, {file_skip} skipped, {file_fail} failed, {file_denied} denied"
        dir_summary = f"Directories: {dir_success} created, {dir_fail} failed, {dir_denied} denied"
        summary = f"\n✨ Create complete:\n   {file_summary}\n   {dir_summary}"
        return "\n".join(results) + summary

    @mcp.tool(name="MyPC-move_file")
    def move_file(source: str, destination: str) -> str:
        """
        Move or rename a file/directory. (WRITE - both paths must be in safe zones)

        Args:
            source: Source path (must be in safe zone).
            destination: Destination path (must be in safe zone).

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(source):
            return f"Error: Cannot move from outside safe zone.\nSource: {source}\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not is_in_safe_zone(destination):
            return f"Error: Cannot move to outside safe zone.\nDestination: {destination}\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(source):
            return f"Error: Source does not exist: {source}"

        try:
            shutil.move(source, destination)
            return f"Moved: {source} -> {destination}"

        except Exception as e:
            return f"Error moving file: {str(e)}"

    @mcp.tool(name="MyPC-delete_file")
    def delete_file(path, to_recycle: bool = True) -> str:
        """
        Delete file(s) or directory/directories. (WRITE - safe zones only)

        Supports both single and multiple paths:
        - Single: path="D:\\ALICE\\file.txt"
        - Multiple: path=["D:\\ALICE\\file1.txt", "D:\\ALICE\\file2.txt"]

        Args:
            path: Path (string) or list of paths to delete (all must be in safe zones).
            to_recycle: If True, move to recycle bin. If False, permanently delete.

        Returns:
            Success or error message.
        """
        # Handle single path (string)
        if isinstance(path, str):
            if not is_in_safe_zone(path):
                return f"Error: Delete operation denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

            if not os.path.exists(path):
                return f"Error: Path does not exist: {path}"

            try:
                if to_recycle:
                    try:
                        from send2trash import send2trash
                        send2trash(path)
                        return f"Moved to Recycle Bin: {path}"
                    except ImportError:
                        return "Error: send2trash not installed. Use to_recycle=False for permanent deletion."
                else:
                    if os.path.isdir(path):
                        shutil.rmtree(path)
                    else:
                        os.remove(path)
                    return f"Permanently deleted: {path}"

            except Exception as e:
                return f"Error deleting: {str(e)}"

        # Handle multiple paths (list)
        elif isinstance(path, list):
            if not path:
                return "Error: Path list is empty"

            results = []
            success_count = 0
            fail_count = 0
            denied_count = 0

            for p in path:
                # Check safe zone
                if not is_in_safe_zone(p):
                    results.append(f"[DENIED] Not in safe zone: {p}")
                    denied_count += 1
                    continue

                if not os.path.exists(p):
                    results.append(f"[FAIL] Path does not exist: {p}")
                    fail_count += 1
                    continue

                try:
                    if to_recycle:
                        try:
                            from send2trash import send2trash
                            send2trash(p)
                            results.append(f"[OK] Moved to Recycle Bin: {os.path.basename(p)}")
                            success_count += 1
                        except ImportError:
                            return "Error: send2trash not installed. Use to_recycle=False for permanent deletion."
                    else:
                        if os.path.isdir(p):
                            shutil.rmtree(p)
                        else:
                            os.remove(p)
                        results.append(f"[OK] Permanently deleted: {os.path.basename(p)}")
                        success_count += 1
                except Exception as e:
                    results.append(f"[FAIL] {os.path.basename(p)}: {str(e)}")
                    fail_count += 1

            summary = f"\nBatch delete complete: {success_count} succeeded, {fail_count} failed, {denied_count} denied"
            return "\n".join(results) + summary

        else:
            return "Error: Path must be a string or list of strings"

    # ==================== COPY OPERATION (Special: INTO safe zone only) ====================

    @mcp.tool(name="MyPC-copy_file")
    def copy_file(source, destination: str) -> str:
        """
        Copy file(s) or directory/directories. (SPECIAL - destination must be in safe zone)

        Security: You can copy FROM anywhere INTO a safe zone, but NOT out of safe zones.

        Supports both single and multiple sources:
        - Single: source="C:\\file.txt", destination="D:\\ALICE\\file.txt"
        - Multiple: source=["C:\\file1.txt", "C:\\file2.txt"], destination="D:\\ALICE\\"
          (copies all files to the destination directory)

        Args:
            source: Source path (string) or list of source paths.
            destination: Destination path. For single source can be file or directory.
                        For multiple sources must be a directory (must be in a safe zone).

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(destination):
            return f"Error: Copy destination must be in a safe zone.\nDestination: {destination}\n\nSafe Zones:\n{get_safe_zones_str()}"

        # Handle single source (string)
        if isinstance(source, str):
            if not os.path.exists(source):
                return f"Error: Source does not exist: {source}"

            try:
                # Ensure parent directory exists
                parent = os.path.dirname(destination)
                if parent and not os.path.exists(parent):
                    os.makedirs(parent, exist_ok=True)

                if os.path.isdir(source):
                    shutil.copytree(source, destination)
                else:
                    shutil.copy2(source, destination)

                return f"Copied: {source} -> {destination}"

            except Exception as e:
                return f"Error copying: {str(e)}"

        # Handle multiple sources (list)
        elif isinstance(source, list):
            if not source:
                return "Error: Source list is empty"

            # Create destination directory if it doesn't exist
            if not os.path.exists(destination):
                try:
                    os.makedirs(destination, exist_ok=True)
                except Exception as e:
                    return f"Error creating destination directory: {str(e)}"

            results = []
            success_count = 0
            fail_count = 0

            for src in source:
                if not os.path.exists(src):
                    results.append(f"[FAIL] Source does not exist: {src}")
                    fail_count += 1
                    continue

                try:
                    # Determine destination path
                    dest_path = os.path.join(destination, os.path.basename(src))

                    if os.path.isdir(src):
                        shutil.copytree(src, dest_path)
                    else:
                        shutil.copy2(src, dest_path)

                    results.append(f"[OK] {os.path.basename(src)}")
                    success_count += 1
                except Exception as e:
                    results.append(f"[FAIL] {os.path.basename(src)}: {str(e)}")
                    fail_count += 1

            summary = f"\nBatch copy complete: {success_count} succeeded, {fail_count} failed"
            return "\n".join(results) + summary

    # ==================== SEARCH OPERATIONS (Anywhere) ====================

    @mcp.tool(name="MyPC-grep_files")
    def grep_files(path: str, keyword: str, recursive: bool = True, extensions: list = None, case_sensitive: bool = False, max_results: int = 100) -> str:
        """
        在指定目录或文件搜索包含特定关键词的内容。 (READ - allowed anywhere)

        Search for keyword content in specified directory or file.

        Args:
            path: 文件或文件夹的绝对路径 | Absolute path to file or folder
            keyword: 要搜索的关键词 | Keyword to search for
            recursive: 如果是文件夹，是否递归搜索子文件夹 | Recursively search subfolders (default True)
            extensions: 限制搜索的文件后缀 | Limit file extensions, e.g. [".txt", ".py"] (default None)
            case_sensitive: 是否区分大小写 | Case sensitive search (default False)
            max_results: 最大返回结果数 | Maximum results to return (default 100)

        Returns:
            搜索结果，包含文件名、行号和匹配行内容 | Search results with filename, line number, and matched content.
        """
        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        import re

        # Prepare search pattern
        flags = 0 if case_sensitive else re.IGNORECASE
        try:
            pattern = re.compile(keyword, flags)
        except Exception as e:
            return f"Error: Invalid search pattern/keyword: {str(e)}"

        results = []
        files_to_scan = []

        # 1. Collect files to scan
        if os.path.isfile(path):
            files_to_scan.append(path)
        else:
            if recursive:
                for root, _, files in os.walk(path):
                    for name in files:
                        files_to_scan.append(os.path.join(root, name))
            else:
                for name in os.listdir(path):
                    full_p = os.path.join(path, name)
                    if os.path.isfile(full_p):
                        files_to_scan.append(full_p)

        # 2. Filter by extensions if provided
        if extensions:
            # Normalize extensions to start with dot
            exts = [e if e.startswith('.') else '.' + e for e in extensions]
            files_to_scan = [f for f in files_to_scan if os.path.splitext(f)[1].lower() in exts]

        # 3. Perform search
        match_count = 0
        for file_path in files_to_scan:
            if match_count >= max_results:
                break

            try:
                # Try reading as text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    for line_num, line in enumerate(f, 1):
                        if pattern.search(line):
                            match_count += 1
                            clean_line = line.strip()
                            rel_path = os.path.relpath(file_path, path) if os.path.isdir(path) else os.path.basename(file_path)
                            results.append(f"[{rel_path}:{line_num}] {clean_line}")

                            if match_count >= max_results:
                                results.append(f"\n... (truncated at {max_results} results)")
                                break
            except Exception:
                # Skip files that can't be read (binary, permissions, etc.)
                continue

        if not results:
            return f"No matches found for '{keyword}' in {path}."

        return f"Found {match_count} matches in {path}:\n\n" + "\n".join(results)
